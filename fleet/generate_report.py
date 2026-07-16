import pandas as pd
import numpy as np
from pptx import Presentation
import os
import json
import shutil
import argparse
import glob

def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument('--no-download', action='store_true', help='Skip downloading new data via Playwright')
    return parser.parse_args()

def main():
    args = parse_args()
    
    # Paths
    base_dir = os.path.dirname(__file__) # This is the fleet/ folder
    workspace_dir = os.path.dirname(base_dir)
    analys_dir = os.path.join(base_dir, 'analys')
    debug_dir = os.path.join(workspace_dir, 'debug_tools')
    inputs_dir = os.path.join(base_dir, 'inputs')
    os.makedirs(debug_dir, exist_ok=True)
    
    # Find Excel & PPTX files
    xlsx_files = [f for f in os.listdir(analys_dir) if f.endswith('.xlsx')]
    pptx_files = [f for f in os.listdir(analys_dir) if f.endswith('.pptx')]
    
    if not xlsx_files or not pptx_files:
        print("Error: Missing Excel or PPTX files in fleet/analys/")
        return
        
    excel_orig = os.path.join(analys_dir, xlsx_files[0])
    pptx_orig = os.path.join(analys_dir, pptx_files[0])
    
    excel_temp = os.path.join(debug_dir, 'temp_report.xlsx')
    pptx_temp = os.path.join(debug_dir, 'temp_presentation.pptx')
    
    import subprocess
    try:
        # Use powershell to copy to bypass file locks on Windows
        subprocess.run(["powershell", "-Command", f'Copy-Item -Path "{excel_orig}" -Destination "{excel_temp}" -Force'], capture_output=True)
        subprocess.run(["powershell", "-Command", f'Copy-Item -Path "{pptx_orig}" -Destination "{pptx_temp}" -Force'], capture_output=True)
        print("Files copied successfully via PowerShell")
    except Exception as e:
        print(f"Warning during file copy: {e}")
        excel_temp = excel_orig
        pptx_temp = pptx_orig
        
    # 1. Extract monthly history data
    print("Extracting history data...")
    df = pd.read_excel(excel_temp, sheet_name='Общее', header=None)
    
    pptx_rentability = {}
    try:
        prs = Presentation(pptx_temp)
        rent_slide = None
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and "рентабельность" in shape.text.lower():
                    if any(s.has_chart for s in slide.shapes):
                        rent_slide = slide
                        break
            if rent_slide:
                break
                
        if rent_slide:
            chart_shape = [s for s in rent_slide.shapes if s.has_chart][0]
            chart = chart_shape.chart
            categories = list(chart.plots[0].categories)
            for series in chart.series:
                series_name = series.name
                month = None
                if 'фев' in series_name: month = '2026-02'
                elif 'мар' in series_name: month = '2026-03'
                elif 'апр' in series_name: month = '2026-04'
                elif 'май' in series_name: month = '2026-05'
                
                if month:
                    for cat_name, val in zip(categories, series.values):
                        if cat_name not in pptx_rentability:
                            pptx_rentability[cat_name] = {}
                        pptx_rentability[cat_name][month] = float(val) if val is not None else None
    except Exception as e:
        print("Error extracting PPTX rentability:", e)
        
    cities = ['Омск', 'Рязань', 'Ижевск', 'Ульяновск', 'Магнитогорск', 'Сургут', 'Киров', 'Чебоксары', 'Орёл']
    
    # 1. Динамическое определение месяцев, за которые есть фактические данные в Excel (по городам)
    dates_row = df.iloc[1, 1:].values
    actual_months = []
    for col_idx, dt in enumerate(dates_row):
        if pd.isna(dt): continue
        dt_str = dt.strftime('%Y-%m') if not isinstance(dt, str) else dt[:7]
        if dt_str.startswith('2025') or dt_str.startswith('2026'):
            # Проверяем, есть ли числовое значение хотя бы для одного города в этой колонке
            has_data = False
            for r in range(2, 11): # Строки городов Омск-Орёл в pandas
                val = df.iloc[r, col_idx + 1]
                if pd.notna(val) and val != '' and not str(val).startswith('='):
                    has_data = True
                    break
            if has_data:
                actual_months.append(dt_str)
                
    print(f"Dynamically discovered actual months: {actual_months}")
    
    def extract_section(start_row, num_rows, date_row):
        section_data = {}
        dates = df.iloc[date_row, 1:].values
        for r in range(start_row, start_row + num_rows):
            city_name = str(df.iloc[r, 0]).strip()
            if not city_name or city_name == 'nan': continue
            section_data[city_name] = {}
            row_vals = df.iloc[r, 1:].values
            for dt, val in zip(dates, row_vals):
                if pd.isna(dt): continue
                dt_str = dt.strftime('%Y-%m') if not isinstance(dt, str) else dt[:7]
                if dt_str in actual_months:
                    if pd.notna(val) and val != '' and not str(val).startswith('='):
                        try:
                            section_data[city_name][dt_str] = float(val)
                        except ValueError:
                            pass
        return section_data
        
    revenue_data = extract_section(2, 10, 1)
    stations_data = extract_section(18, 10, 17)
    rps_data = extract_section(57, 10, 56)
    rent_data = extract_section(87, 9, 86)
    
    for city, months in pptx_rentability.items():
        if city in rent_data:
            for m, val in months.items():
                if m not in rent_data[city] or pd.isna(rent_data[city][m]):
                    rent_data[city][m] = val
                    
    # 2. Extract Jewelry distribution data
    print("Extracting jewelry distribution data...")
    jewelry_data = {}
    all_dfs = []
    
    for city in cities:
        pattern = os.path.join(inputs_dir, f"revenue_{city}_*.csv")
        files = glob.glob(pattern)
        if not files:
            print(f"No CSV found for {city} in inputs/")
            continue
        files.sort()
        latest_file = files[-1]
        
        try:
            csv_df = pd.read_csv(latest_file)
            csv_df = csv_df[csv_df['office_status'].astype(str).str.strip().str.lower() == 'placed']
            csv_df = csv_df[csv_df['remove_date'].astype(str).str.strip().str.contains('2222-02-01|01.02.2222', regex=True)]
            csv_df['jewelry'] = csv_df['jewelry'].fillna('0').astype(str).str.strip().str.lower()
            csv_df['jewelry'] = csv_df['jewelry'].replace(['не указано', '', 'nan'], '0')
            
            all_dfs.append(csv_df[['jewelry']])
            counts = csv_df['jewelry'].value_counts().to_dict()
            jewelry_data[city] = counts
        except Exception as e:
            print(f"Error parsing jewelry for {city}: {e}")
            
    if all_dfs:
        df_overall = pd.concat(all_dfs, ignore_index=True)
        jewelry_data['Общее'] = df_overall['jewelry'].value_counts().to_dict()
        
    # Compile final structure
    final_data = {}
    months = actual_months
    
    for city in cities + ['Общее']:
        final_data[city] = {
            'history': [],
            'jewelry': jewelry_data.get(city, {})
        }
        for m in months:
            if city == 'Общее':
                # Выручка: сумма всех городов
                rev_list = [revenue_data.get(c, {}).get(m, None) for c in cities]
                rev_clean = [v for v in rev_list if v is not None]
                rev = sum(rev_clean) if rev_clean else None
                
                # Кол-во станций: сумма всех городов
                st_list = [stations_data.get(c, {}).get(m, None) for c in cities]
                st_clean = [v for v in st_list if v is not None]
                st = sum(st_clean) if st_clean else None
                
                # RPS: среднее по городам (для соответствия расчетной логике Excel-отчета)
                rps_list = [rps_data.get(c, {}).get(m, None) for c in cities]
                rps_clean = [v for v in rps_list if v is not None]
                rps = sum(rps_clean) / len(rps_clean) if rps_clean else None
                
                # Рентабельность: среднее по городам
                rent_list = [rent_data.get(c, {}).get(m, None) for c in cities]
                rent_clean = [v for v in rent_list if v is not None]
                rent = sum(rent_clean) / len(rent_clean) if rent_clean else None
            else:
                rev = revenue_data.get(city, {}).get(m, None)
                st = stations_data.get(city, {}).get(m, None)
                rps = rps_data.get(city, {}).get(m, None)
                rent = rent_data.get(city, {}).get(m, None)
                
            if rps is None and rev is not None and st is not None and st > 0:
                rps = round(rev / st, 2)
                
            if any(v is not None for v in [rev, st, rps, rent]):
                final_data[city]['history'].append({
                    'month': m,
                    'revenue': rev,
                    'stations': st,
                    'revenue_per_station': rps,
                    'rentability': rent
                })
                
    render_html(final_data)

def render_html(data):
    base_dir = os.path.dirname(__file__)
    template_path = os.path.join(base_dir, 'templates', 'dashboard_template.html')
    if not os.path.exists(template_path):
        print("Template not found!")
        return
        
    with open(template_path, 'r', encoding='utf-8') as f:
        template = f.read()
        
    json_data = json.dumps(data, ensure_ascii=False)
    html = template.replace('{{ REPORT_DATA }}', json_data)
    
    output_path = os.path.join(base_dir, 'report.html')
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html)
    print(f"Report successfully generated at {output_path}")

if __name__ == '__main__':
    main()
